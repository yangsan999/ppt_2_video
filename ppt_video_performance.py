import configparser
import json
import os
import uuid
import asyncio
import aiohttp
import shutil


from subprocess import call
from pptx import Presentation  # pip install python-pptx
from pptx_tools import utils  # pip install python-pptx-interface

# pip install tencentcloud-sdk-python
from tencentcloud.common import credential
from tencentcloud.common.exception.tencent_cloud_sdk_exception import TencentCloudSDKException
from tencentcloud.tts.v20190823 import models, tts_client


# Read configuration from env.ini
config = configparser.ConfigParser()
config.read('env.ini')

SECRET_ID = config['Tencent_API']['SecretId']
SECRET_KEY = config['Tencent_API']['SecretKey']

# The maximum number of concurrent HTTP requests
MAX_CONCURRENT_REQUESTS = 20

# The semaphore to limit the number of concurrent HTTP requests
semaphore = asyncio.Semaphore(MAX_CONCURRENT_REQUESTS)

def set_temp_dir():
    # Set the temp directory
    temp_path = os.path.join(os.getcwd(), 'temp')
    image_folder = os.path.join(temp_path, 'images')
    if os.path.exists(temp_path):
        shutil.rmtree(temp_path)
    os.makedirs(temp_path)
    os.makedirs(image_folder)
    return (temp_path, image_folder)


def ppt_save_as_png(image_folder, input_pptx):

    if utils.save_pptx_as_png(image_folder, input_pptx, overwrite_folder=True):
        print("ppt save as png is success")
    else:
        raise SystemExit

    # Iterate over the files in the directory
    for file in os.listdir(image_folder):
        # Get the file name and extension
        name, ext = os.path.splitext(file)

        # Extract the number from the file name, number -1 to make start with 0
        num = str(int(name.split('幻灯片')[1])-1)

        # Rename the file using the new naming convention
        os.rename(os.path.join(image_folder, file), os.path.join(
            image_folder, 'slide_{}{}'.format(num, ext)))

    print("png rename is success")


async def create_tts_task(note, session_id, session):
    # Acquire the semaphore
    async with semaphore:

        # Create a Tencent Cloud credential object
        cred = credential.Credential(SECRET_ID, SECRET_KEY)

        # Create a TTS client
        client = tts_client.TtsClient(cred, "ap-guangzhou")

        # Set the request parameters
        req = models.CreateTtsTaskRequest()
        params = {
            "Text": note,
            "SessionId": str(session_id),
            "ModelType": 1,
            "Speed": 0,
            "SampleRate": 16000,
            "Codec": "mp3",
            "Volume": 0,
            "VoiceType": 101024,
        }
        req.from_json_string(json.dumps(params))

        # Call the TTS service
        try:
            resp = client.CreateTtsTask(req)
            print(str(session_id)+" TTS Task is success")
            task_id = resp.Data.TaskId
        except TencentCloudSDKException as err:
            print(err)
            return None

        # Wait for the task to complete
        while True:
            task_req = models.CreateTtsTaskResponse()
            task_req.TaskId = task_id
            task_resp = client.DescribeTtsTaskStatus(task_req)
            if task_resp.Data.StatusStr == "success":
                print(str(task_id)+" Describe TTS task status is success")
                return task_resp.Data.ResultUrl
            elif task_resp.Data.StatusStr == "failed":
                print("TTS task failed: {}".format(task_resp.ErrorMessage))
                return None
            # Wait for 1 second before checking the task status again
            await asyncio.sleep(1)


async def get_tts_result(task_id, index, session):
    # Acquire the semaphore
    async with semaphore:
        async with aiohttp.ClientSession() as session:
            async with session.get(task_id) as resp:
                audio_bytes = await resp.read()

    # Save the bytes to a file
    with open(os.path.join(temp_path, 'slide_{}.mp3'.format(index)), 'wb') as f:
        f.write(audio_bytes)
        print('slide_{}.mp3'.format(index) + " write is success")


async def create_images_audios(input_pptx):
    # Convert PowerPoint presentation to images
    ppt_save_as_png(image_folder, input_pptx)

    # Open the PowerPoint presentation
    prs = Presentation(input_pptx)
    # Create a list to store the notes for each slide
    notes_list = []

    # Iterate over the slides in the presentation
    for i, slide in enumerate(prs.slides):
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if not notes.strip():
                # Skip empty notes
                continue
            # add break time at the beginning and end of notes
            notes = r"<speak><break time='1s'/></speak>" + \
                notes + r"<speak><break time='2s'/></speak>"
            notes_list.append(notes)

    # Create a session ID for the TTS tasks
    session_id = uuid.uuid1()

    # Create a list to store the task IDs
    task_ids = []

    # Create a TTS task for each slide
    async with aiohttp.ClientSession() as session:
        tasks = []
        for i, note in enumerate(notes_list):
            task = asyncio.ensure_future(create_tts_task(
                note, session_id, session))
            tasks.append(task)
        task_ids = await asyncio.gather(*tasks)

    # Wait for the TTS tasks to complete
    async with aiohttp.ClientSession() as session:
        tasks = []
        for i, task_id in enumerate(task_ids):
            task = asyncio.ensure_future(get_tts_result(
                task_id, i, session))
            tasks.append(task)
        await asyncio.gather(*tasks)


def create_slide_videos():
    num_videos = sum([len(files)
                     for root, dirs, files in os.walk(image_folder)])
    # Combine each slide image and audio into a video
    print('Start combine images and audios into videos, you may need to wait long, long, long time')
    for i in range(num_videos):
        image_path = os.path.join(image_folder, 'slide_{}.png'.format(i))
        audio_path = os.path.join(temp_path, 'slide_{}.mp3'.format(i))
        video_path = os.path.join(temp_path, 'slide_{}.mp4'.format(i))

        call(["ffmpeg",
              '-hide_banner',
              '-loglevel', 'error',
              '-y',
              '-loop', '1',  '-i', image_path,
              '-i', audio_path,
              '-c:v', 'libx264', '-b:v', '2400k',
              '-tune', 'stillimage', '-s', '1920x1080',
             '-c:a', 'aac', '-b:a', '320k',
              '-pix_fmt', 'yuv420p', '-shortest', '-r', '30',
              '-threads', '4', '-movflags', '+faststart', video_path])

        print(str('slide_{}.mp4'.format(i)) + " is success")

    video_list = [os.path.join(temp_path, 'slide_{}.mp4'.format(i))
                  for i in range(num_videos)]
    video_list_str = "\n".join(["file '{}'".format(f) for f in video_list])
    video_list_str_path = os.path.join(temp_path, 'video_list_str.txt')

    with open(video_list_str_path, 'w') as f:
        # Write the 'file' prefix and the file name to the file, one per line
        f.write(video_list_str)
        print("video_list_str.txt write is success")

    return (video_list_str_path)


def create_final_video(video_list_str_path, output_video):
    print('Start combine final video, you are almost success')
    call(["ffmpeg",
          '-hide_banner',
          '-loglevel', 'error',
          '-y', '-f', 'concat',
         '-safe', '0', '-i', video_list_str_path, '-c', 'copy', output_video])
    print(str(output_video) + " is success")


def main_without_cache():  # Always create new png, audio to create videos
    asyncio.run(create_images_audios(input_pptx))
    video_list_str_path = create_slide_videos()
    create_final_video(video_list_str_path, output_video)

    # Clean up the temporary files
    shutil.rmtree(temp_path)
    print("clean up is success, you have done all job!")


def main_with_cache():    # Use cached png, audio to create videos
    # IF run first time, create temporary dirs, png and audio
    if not os.path.exists(image_folder):
        os.makedirs(temp_path)
        os.makedirs(image_folder)
        asyncio.run(create_images_audios(input_pptx))
    video_list_str_path = create_slide_videos()
    create_final_video(video_list_str_path, output_video)


if __name__ == '__main__':
    # Batch process all pptx
    directory = os.path.join(os.getcwd(), 'input')
    filenames = os.listdir(directory)
    for filename in filenames:
        if filename.endswith(".pptx"):
            input_pptx = os.path.join(directory, filename)
            output_video = os.path.join(
                os.getcwd(), 'output', os.path.splitext(filename)[0] + ".mp4")

            # Process with main_without_cache
            temp_path, image_folder = set_temp_dir()
            main_without_cache()

    # Process one by one with main_with_cache
    # input_pptx = os.path.join(os.getcwd(), 'input', 'example.pptx')
    # output_video = os.path.join(os.getcwd(), 'output', 'example.mp4')
    # temp_path = os.path.join(os.getcwd(), 'temp')
    # image_folder = os.path.join(temp_path, 'images')
    # main_with_cache()
