import uuid
import asyncio
import aiohttp

from pptx import Presentation  # pip install python-pptx

from ppt_save_as_png import ppt_save_as_png
from tencent_tts import create_tts_task, get_tts_result


async def create_images_audios(temp_path, image_folder, input_pptx):
    # Convert PowerPoint presentation to images
    ppt_save_as_png(image_folder, input_pptx)

    # Open the PowerPoint presentation
    prs = Presentation(input_pptx)
    # Create a list to store the notes for each slide
    notes_list = []

    # Iterate over the slides in the presentation
    for i, slide in enumerate(prs.slides):
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if not notes.strip():
                # Skip empty notes
                continue
            # add break time at the beginning and end of notes
            notes = r"<speak><break time='1s'/></speak>" + \
                notes + r"<speak><break time='2s'/></speak>"
            notes_list.append(notes)

    # Create a session ID for the TTS tasks
    session_id = uuid.uuid1()

    # Create a list to store the task IDs
    task_ids = []

    # Create a TTS task for each slide
    async with aiohttp.ClientSession() as session:
        tasks = []
        for i, note in enumerate(notes_list):
            task = asyncio.ensure_future(create_tts_task(
                note, session_id, session))
            tasks.append(task)
        task_ids = await asyncio.gather(*tasks)

    # Wait for the TTS tasks to complete
    async with aiohttp.ClientSession() as session:
        tasks = []
        for i, task_id in enumerate(task_ids):
            task = asyncio.ensure_future(
                get_tts_result(task_id, i, session, temp_path))
            tasks.append(task)
        await asyncio.gather(*tasks)
